"""
PowerPoint Generation Module using python-pptx

This module provides a PPTGenerator class with all the basic functions
for creating, modifying, and managing PowerPoint presentations.

Required package: python-pptx
Install with: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from typing import Optional, Tuple, List, Union
from config import get_config, get_logger


class PPTGenerator:
    """
    A comprehensive PowerPoint generator class that provides methods for
    creating and manipulating PowerPoint presentations.
    """
    
    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize the PPT generator with an optional template.
        
        Args:
            template_path (str, optional): Path to a PowerPoint template file.
                                         If None, creates a blank presentation.
        """
        self.config = get_config()
        self.logger = get_logger()
        
        if template_path and os.path.exists(template_path):
            self.presentation = Presentation(template_path)
            self.logger.info(f"Loaded presentation template: {template_path}")
        else:
            self.presentation = Presentation()
            self.logger.info("Created new blank presentation")
        
        # Store the original slide layouts for reference
        self.slide_layouts = self.presentation.slide_layouts
        
    def save(self, filename: str) -> None:
        """
        Save the presentation to a file.
        
        Args:
            filename (str): The filename to save the presentation as.
        """
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        self.presentation.save(filename)
        self.logger.info(f"Presentation saved as: {filename}")
        
    def get_slide_count(self) -> int:
        """Get the total number of slides in the presentation."""
        return len(self.presentation.slides)
        
    def get_available_layouts(self) -> List[str]:
        """Get a list of available slide layout names."""
        layout_names = []
        for i, layout in enumerate(self.slide_layouts):
            layout_names.append(f"{i}: {layout.name}")
        return layout_names
    
    def add_slide(self, layout_index: int = 0, title: str = "", subtitle: str = "") -> int:
        """
        Add a new slide to the presentation.
        
        Args:
            layout_index (int): Index of the slide layout to use (default: 0 for title slide)
            title (str): Title text for the slide
            subtitle (str): Subtitle text for the slide
            
        Returns:
            int: Index of the newly created slide
            
        Available layout indices:
            0: Title Slide
            1: Title and Content
            2: Section Header
            3: Two Content
            4: Comparison
            5: Title Only
            6: Blank
            7: Content with Caption
            8: Picture with Caption
        """
        try:
            # Get the layout
            if layout_index >= len(self.slide_layouts):
                layout_index = 1  # Default to Title and Content
                
            slide_layout = self.slide_layouts[layout_index]
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Add title if provided
            if title and slide.shapes.title:
                slide.shapes.title.text = title
                
            # Add subtitle/content if provided
            if subtitle:
                # Try to find placeholders for subtitle/content
                placeholders = [shape for shape in slide.placeholders]
                if len(placeholders) > 1:  # Has subtitle/content placeholder
                    placeholders[1].text = subtitle
                    
            slide_index = len(self.presentation.slides) - 1
            self.logger.info(f"Added slide {slide_index + 1} with layout: {slide_layout.name}")
            return slide_index
            
        except Exception as e:
            self.logger.error(f"Error adding slide: {str(e)}")
            return -1
    
    def delete_slide(self, n: int) -> bool:
        """
        Delete a slide from the presentation by index.
        
        Args:
            n (int): Index of the slide to delete (0-based indexing)
            
        Returns:
            bool: True if slide was successfully deleted, False otherwise
        """
        try:
            slides = self.presentation.slides
            total_slides = len(slides)
            
            # Validate slide index
            if n < 0 or n >= total_slides:
                self.logger.error(f"Slide index {n} is out of range (0-{total_slides-1})")
                return False
                
            if total_slides == 1:
                self.logger.error("Cannot delete the last remaining slide")
                return False
                
            # Get slide XML element and delete it
            slide_to_delete = slides[n]
            slide_id = slide_to_delete.slide_id
            
            # Remove from slides collection
            slides._sldIdLst.remove(slides._sldIdLst[n])
            
            self.logger.info(f"Successfully deleted slide {n + 1}")
            return True
            
        except Exception as e:
            self.logger.error(f"Error deleting slide {n}: {str(e)}")
            return False
    
    def modify_layout(self, slide_index: int, new_layout_index: int) -> bool:
        """
        Change the layout of an existing slide.
        
        Args:
            slide_index (int): Index of the slide to modify (0-based)
            new_layout_index (int): Index of the new layout to apply
            
        Returns:
            bool: True if layout was successfully changed, False otherwise
        """
        try:
            slides = self.presentation.slides
            
            # Validate slide index
            if slide_index < 0 or slide_index >= len(slides):
                print(f"Error: Slide index {slide_index} is out of range")
                return False
                
            # Validate layout index
            if new_layout_index < 0 or new_layout_index >= len(self.slide_layouts):
                print(f"Error: Layout index {new_layout_index} is out of range")
                return False
                
            slide = slides[slide_index]
            new_layout = self.slide_layouts[new_layout_index]
            
            # Store existing content before changing layout
            existing_title = ""
            existing_content = []
            
            # Try to preserve title
            if slide.shapes.title:
                existing_title = slide.shapes.title.text
                
            # Try to preserve other text content
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    if shape != slide.shapes.title and shape.text_frame.text.strip():
                        existing_content.append(shape.text_frame.text)
            
            # Change the layout
            slide.follow_master_slide = False
            slide.slide_layout = new_layout
            
            # Restore content where possible
            if existing_title and slide.shapes.title:
                slide.shapes.title.text = existing_title
                
            # Restore other content to available placeholders
            placeholders = [shape for shape in slide.placeholders 
                          if shape != slide.shapes.title and hasattr(shape, "text_frame")]
            
            for i, content in enumerate(existing_content):
                if i < len(placeholders):
                    placeholders[i].text = content
                    
            print(f"Successfully changed slide {slide_index + 1} to layout: {new_layout.name}")
            return True
            
        except Exception as e:
            print(f"Error modifying layout for slide {slide_index}: {str(e)}")
            return False
    
    def insert_chart(self, slide_index: int, chart_type: str = "column", 
                    data: dict = None, position: Tuple[float, float] = None, 
                    size: Tuple[float, float] = None) -> bool:
        """
        Insert a chart into a slide.
        
        Args:
            slide_index (int): Index of the slide to add chart to
            chart_type (str): Type of chart ('column', 'bar', 'line', 'pie')
            data (dict): Chart data in format {'categories': [...], 'series': [{'name': '...', 'values': [...]}]}
            position (tuple): (x, y) position in inches, default (1, 2)
            size (tuple): (width, height) in inches, default (8, 5)
            
        Returns:
            bool: True if chart was successfully added, False otherwise
        """
        try:
            slides = self.presentation.slides
            
            # Validate slide index
            if slide_index < 0 or slide_index >= len(slides):
                print(f"Error: Slide index {slide_index} is out of range")
                return False
                
            slide = slides[slide_index]
            
            # Default data if none provided
            if data is None:
                data = {
                    'categories': ['Q1', 'Q2', 'Q3', 'Q4'],
                    'series': [
                        {'name': 'Sales', 'values': [10, 15, 12, 18]},
                        {'name': 'Profit', 'values': [5, 8, 6, 10]}
                    ]
                }
            
            # Default position and size
            if position is None:
                position = (Inches(1), Inches(2))
            else:
                position = (Inches(position[0]), Inches(position[1]))
                
            if size is None:
                size = (Inches(8), Inches(5))
            else:
                size = (Inches(size[0]), Inches(size[1]))
            
            # Map chart type string to enum
            chart_type_map = {
                'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'bar': XL_CHART_TYPE.BAR_CLUSTERED,
                'line': XL_CHART_TYPE.LINE,
                'pie': XL_CHART_TYPE.PIE
            }
            
            chart_type_enum = chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            # Create chart data
            chart_data = CategoryChartData()
            chart_data.categories = data['categories']
            
            for series in data['series']:
                chart_data.add_series(series['name'], series['values'])
            
            # Add chart to slide
            chart = slide.shapes.add_chart(
                chart_type_enum, position[0], position[1], size[0], size[1], chart_data
            ).chart
            
            # Customize chart appearance
            chart.has_legend = True
            chart.legend.include_in_layout = False
            
            print(f"Successfully added {chart_type} chart to slide {slide_index + 1}")
            return True
            
        except Exception as e:
            print(f"Error inserting chart to slide {slide_index}: {str(e)}")
            return False
    
    def insert_image(self, slide_index: int, image_path: str, 
                    position: Tuple[float, float] = None, 
                    size: Tuple[float, float] = None) -> bool:
        """
        Insert an image into a slide.
        
        Args:
            slide_index (int): Index of the slide to add image to
            image_path (str): Path to the image file
            position (tuple): (x, y) position in inches, default (1, 1)
            size (tuple): (width, height) in inches, default (auto-sized)
            
        Returns:
            bool: True if image was successfully added, False otherwise
        """
        try:
            slides = self.presentation.slides
            
            # Validate slide index
            if slide_index < 0 or slide_index >= len(slides):
                print(f"Error: Slide index {slide_index} is out of range")
                return False
                
            # Validate image path
            if not os.path.exists(image_path):
                print(f"Error: Image file not found: {image_path}")
                return False
                
            slide = slides[slide_index]
            
            # Default position
            if position is None:
                position = (Inches(1), Inches(1))
            else:
                position = (Inches(position[0]), Inches(position[1]))
            
            # Add image to slide
            if size is None:
                # Auto-size the image
                pic = slide.shapes.add_picture(image_path, position[0], position[1])
            else:
                # Use specified size
                pic = slide.shapes.add_picture(
                    image_path, position[0], position[1], 
                    Inches(size[0]), Inches(size[1])
                )
            
            print(f"Successfully added image to slide {slide_index + 1}: {os.path.basename(image_path)}")
            return True
            
        except Exception as e:
            print(f"Error inserting image to slide {slide_index}: {str(e)}")
            return False
    
    def update_text(self, slide_index: int, text_updates: dict, 
                   font_size: Optional[int] = None, 
                   font_color: Optional[Tuple[int, int, int]] = None,
                   bold: Optional[bool] = None,
                   alignment: Optional[str] = None) -> bool:
        """
        Update text content in a slide.
        
        Args:
            slide_index (int): Index of the slide to update
            text_updates (dict): Dictionary with keys like 'title', 'content', or shape indices
                                Example: {'title': 'New Title', 'content': 'New content', 0: 'Text for shape 0'}
            font_size (int, optional): Font size in points
            font_color (tuple, optional): RGB color tuple (r, g, b)
            bold (bool, optional): Make text bold
            alignment (str, optional): Text alignment ('left', 'center', 'right')
            
        Returns:
            bool: True if text was successfully updated, False otherwise
        """
        try:
            slides = self.presentation.slides
            
            # Validate slide index
            if slide_index < 0 or slide_index >= len(slides):
                print(f"Error: Slide index {slide_index} is out of range")
                return False
                
            slide = slides[slide_index]
            
            # Update title if specified
            if 'title' in text_updates and slide.shapes.title:
                slide.shapes.title.text = text_updates['title']
                self._apply_text_formatting(slide.shapes.title.text_frame, 
                                          font_size, font_color, bold, alignment)
            
            # Update content/subtitle if specified
            if 'content' in text_updates or 'subtitle' in text_updates:
                content_text = text_updates.get('content', text_updates.get('subtitle', ''))
                # Find content placeholders
                for shape in slide.placeholders:
                    if (shape != slide.shapes.title and hasattr(shape, 'text_frame') 
                        and shape.text_frame):
                        shape.text = content_text
                        self._apply_text_formatting(shape.text_frame, 
                                                  font_size, font_color, bold, alignment)
                        break
            
            # Update specific shapes by index
            shapes_list = list(slide.shapes)
            for key, value in text_updates.items():
                if isinstance(key, int) and 0 <= key < len(shapes_list):
                    shape = shapes_list[key]
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        shape.text_frame.text = str(value)
                        self._apply_text_formatting(shape.text_frame, 
                                                  font_size, font_color, bold, alignment)
                    elif hasattr(shape, 'text'):
                        shape.text = str(value)
            
            print(f"Successfully updated text in slide {slide_index + 1}")
            return True
            
        except Exception as e:
            print(f"Error updating text in slide {slide_index}: {str(e)}")
            return False
    
    def _apply_text_formatting(self, text_frame, font_size=None, font_color=None, 
                             bold=None, alignment=None):
        """Helper method to apply text formatting."""
        try:
            if not text_frame or not text_frame.paragraphs:
                return
                
            for paragraph in text_frame.paragraphs:
                if alignment:
                    align_map = {
                        'left': PP_ALIGN.LEFT,
                        'center': PP_ALIGN.CENTER,
                        'right': PP_ALIGN.RIGHT
                    }
                    if alignment.lower() in align_map:
                        paragraph.alignment = align_map[alignment.lower()]
                
                for run in paragraph.runs:
                    if font_size:
                        run.font.size = Pt(font_size)
                    if font_color:
                        run.font.color.rgb = RGBColor(*font_color)
                    if bold is not None:
                        run.font.bold = bold
                        
        except Exception as e:
            print(f"Warning: Could not apply formatting: {str(e)}")
    
    def change_background(self, slide_index: int, background_type: str = "solid", 
                         color: Tuple[int, int, int] = (255, 255, 255),
                         image_path: str = None) -> bool:
        """
        Change the background of a slide.
        
        Args:
            slide_index (int): Index of the slide to change background
            background_type (str): Type of background ('solid', 'image')
            color (tuple): RGB color tuple for solid backgrounds, default white (255, 255, 255)
            image_path (str): Path to background image (for 'image' type)
            
        Returns:
            bool: True if background was successfully changed, False otherwise
        """
        try:
            slides = self.presentation.slides
            
            # Validate slide index
            if slide_index < 0 or slide_index >= len(slides):
                print(f"Error: Slide index {slide_index} is out of range")
                return False
                
            slide = slides[slide_index]
            background = slide.background
            
            if background_type.lower() == "solid":
                # Set solid color background
                background.fill.solid()
                background.fill.fore_color.rgb = RGBColor(*color)
                print(f"Applied solid background color RGB{color} to slide {slide_index + 1}")
                
            elif background_type.lower() == "image":
                # Set image background
                if not image_path or not os.path.exists(image_path):
                    print(f"Error: Image file not found: {image_path}")
                    return False
                    
                background.fill.patterned()
                # Note: python-pptx has limitations with image backgrounds
                # Alternative approach: add image as a full-slide shape
                slide_width = self.presentation.slide_width
                slide_height = self.presentation.slide_height
                
                # Add image as background (behind all other shapes)
                img_shape = slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)
                
                # Move image to back
                shape_list = list(slide.shapes)
                if len(shape_list) > 1:
                    # Move the image to the beginning (back)
                    slide.shapes._spTree.remove(img_shape._element)
                    slide.shapes._spTree.insert(2, img_shape._element)  # Insert after background elements
                
                print(f"Applied image background to slide {slide_index + 1}: {os.path.basename(image_path)}")
                
            else:
                print(f"Error: Unsupported background type: {background_type}")
                return False
                
            return True
            
        except Exception as e:
            print(f"Error changing background of slide {slide_index}: {str(e)}")
            return False